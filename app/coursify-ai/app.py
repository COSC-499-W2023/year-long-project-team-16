import base64
import glob
import io
import logging
from mailbox import Message
from pydoc_data import topics
import subprocess
import time
from urllib.parse import unquote_to_bytes
from flask import after_this_request, flash, session
from datetime import datetime
import os
import re
import random
import string
from click import wrap_text
from flask import Flask, jsonify, render_template, request, send_from_directory, url_for, Response, send_file, make_response, redirect
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.pdfbase.pdfmetrics import stringWidth
from flask_cors import CORS
import openai 
import matplotlib.pyplot as plt
from io import BytesIO
from pymongo import MongoClient
from gridfs import GridFS, NoFile
from PyPDF2 import PdfReader
from werkzeug.utils import secure_filename
from bson.objectid import ObjectId

from flask_bcrypt import Bcrypt
from flask_login import LoginManager, UserMixin, login_user, login_required, logout_user, current_user
from pymongo import MongoClient
from bson.objectid import ObjectId
from itsdangerous import URLSafeTimedSerializer, SignatureExpired, BadSignature
from sendgrid import SendGridAPIClient
from sendgrid.helpers.mail import Mail
from flask_mail import Mail, Message
from pptx import Presentation
from flask import send_file, abort
from flask_wtf import FlaskForm
from wtforms import StringField, PasswordField, SubmitField, SelectField, RadioField, TextAreaField
from wtforms.validators import DataRequired, Email, EqualTo, Length, ValidationError
from datetime import datetime
from docx import Document
import tempfile
import subprocess
import os
from PIL import Image
from urllib.parse import quote
import shutil
from pathlib import Path

import subprocess




class RegistrationForm(FlaskForm):
    '''Registration form class. Inherits from FlaskForm. '''
    first_name = StringField('First Name', validators=[DataRequired()])
    last_name = StringField('Last Name', validators=[DataRequired()])
    email = StringField('Email', validators=[DataRequired(), Email()])
    password = PasswordField('Password', validators=[
        DataRequired(),
        Length(min=8, message='Password must be at least 8 characters long.'),
        EqualTo('confirm_password', message='Passwords must match.')
    ])
    confirm_password = PasswordField('Confirm Password')
    birth_date = SelectField('Birth Date', choices=[(str(i), str(i)) for i in range(1, 32)])
    birth_month = SelectField('Birth Month', choices=[(str(i), str(i)) for i in range(1, 13)])
    birth_year = SelectField('Birth Year', choices=[(str(i), str(i)) for i in range(1900, 2022)])
    gender = RadioField('Gender', choices=[('Male', 'Male'), ('Female', 'Female'), ('Custom', 'Custom')])
    custom_gender = TextAreaField('Custom Gender')
    submit = SubmitField('Register')



app = Flask(__name__, template_folder='my_templates')
CORS(app)
bcrypt = Bcrypt(app)
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'
openai.api_key = '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'


app.secret_key = 'coursifyai1234'

serializer = URLSafeTimedSerializer(app.secret_key)




# Setup MongoDB connection
client = MongoClient('mongodb+srv://Remy:1234@cluster0.vgzdbrr.mongodb.net/')
db = client['new_pdfs']
db2=client['Login_details']
db3=client['reviews']
fs = GridFS(db)
users_collection=db2.users
reviews_collection=db3.reviews

app.config['MAIL_SERVER'] = 'smtp-mail.outlook.com'
app.config['MAIL_PORT'] = 587
app.config['MAIL_USE_TLS'] = True
app.config['MAIL_USERNAME'] = 'coursify@outlook.com'  
app.config['MAIL_PASSWORD'] = 'Gunners4Eva$'

mail = Mail(app)



# User model
class User(UserMixin):
    '''User class for Flask-Login.'''
    def __init__(self, user_id, email):
        self.user_id = str(user_id)
        self.email = email

    def get_id(self):
        return self.user_id

# User Loader
@login_manager.user_loader
def load_user(user_id):
    '''Function to load a user from the database.'''
    user = users_collection.find_one({"_id": ObjectId(user_id)})
    if not user:
        return None
    return User(user_id=user["_id"], email=user["email"])
@app.route('/')
def home():
    '''Function to render the homepage.'''

    # If the user is authenticated, redirect to the dashboard
    if current_user.is_authenticated:
        return redirect(url_for('index'))
    # Otherwise, show the homepage with login and register options
    return render_template('homepage.html')

SENDGRID_API_KEY = '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'  
def validate_password(password):
    '''Validation of password'''
    if len(password) < 8:
        return "Password must be at least 8 characters long."
    if not re.search("[a-z]", password):
        return "Password must contain at least one lowercase letter."
    if not re.search("[A-Z]", password):
        return "Password must contain at least one uppercase letter."
    if not re.search("[0-9]", password):
        return "Password must contain at least one number."
    if not re.search("[!@#$%^&*]", password):
        return "Password must contain at least one special character (!@#$%^&*)."
    return None
   
    # Registration Route
@app.route('/register', methods=['GET', 'POST'])
def register():
    '''Function to register a new user.'''
    if request.method == 'POST':
        first_name = request.form.get('first_name')
        last_name = request.form.get('last_name')
        email = request.form.get('email')
        password = request.form.get('password')
        birth_day = int(request.form.get('birth_day'))
        birth_month = int(request.form.get('birth_month'))
        birth_year = int(request.form.get('birth_year'))
        gender = request.form.get('gender')
        custom_gender = request.form.get('custom_gender') if 'custom_gender' in request.form else None
        if gender == 'custom' and custom_gender:
            gender_value = custom_gender
        else:
            gender_value = gender
        # Validate if user exists
        user = users_collection.find_one({"email": email})
        if user:
            flash('Email already exists.')
            return redirect(url_for('register'))
        
        password = request.form.get('password')
        password_error = validate_password(password)
        if password_error:
            flash(password_error)
            return redirect(url_for('register'))
        hashed_password = bcrypt.generate_password_hash(password).decode('utf-8')

        # Construct birthdate and gender
        birth_date = datetime(birth_year, birth_month, birth_day)
        gender = custom_gender if gender == 'Custom' else gender
     
        # Insert new user into MongoDB
        user_id = users_collection.insert_one({
            "first_name": first_name,
            "last_name": last_name,
            "email": email,
            "password": hashed_password,
            "birth_date": birth_date,
            "gender": gender_value,
            "verified": False  # Assuming you handle email verification separately
        }).inserted_id
        
        # Auto-login after register
        user_obj = User(user_id=str(user_id), email=email)
        login_user(user_obj)
        print(current_user.is_authenticated)
        flash('Registration successful. Welcome!', 'success')
        return redirect(url_for('index'))  # Redirect to index page
    return render_template('register.html')

# Send Email Function
def send_email(to_email, subject, confirm_url):
    '''Function to send an email to the user.'''
    html_content = render_template('email_verification.html', confirm_url=confirm_url)
    message = Mail(
        from_email='',  # Replace with your verified sender email
       to_emails=to_email,
        subject=subject,
        html_content=html_content
    )
    #try:
        #sg = SendGridAPIClient(SENDGRID_API_KEY)
       # response = sg.send(message)
       #print(f"Email sent with status code: {response.status_code}")
    #except Exception as e:
     #   print(f"Error sending email: {e}")

# Email Confirmation Route
@app.route('/confirm/<token>')
def confirm_email(token):
    '''Function to confirm the user's email address.'''
    try:
        email = serializer.loads(
            token, 
            salt='email-confirmation-salt', 
            max_age=3600  # Token expires after 1 hour
        )
    except (SignatureExpired, BadSignature):
        flash('The confirmation link is invalid or has expired.', 'danger')
        return redirect(url_for('index'))

    user = users_collection.find_one({"email": email})
    if user and not user['verified']:
        users_collection.update_one({"_id": user['_id']}, {"$set": {"verified": True}})
        flash('Your account has been activated!', 'success')
    else:
        flash('Account already activated or not found.', 'success')
    return redirect(url_for('login'))

# Login Route
# Updated Login Route with "Remember Me"
@app.route('/login', methods=['GET', 'POST'])
def login():
    '''Function to log in the user.'''
    if request.method == 'POST':
        email = request.form['email']
        password = request.form['password']
        remember = 'remember' in request.form  # Check if 'Remember Me' is checked

        user = users_collection.find_one({"email": email})
        if user and bcrypt.check_password_hash(user['password'], password):
            user_obj = User(user_id=user["_id"], email=email)
            login_user(user_obj, remember=remember)
            return redirect(url_for('index'))

        flash('Invalid credentials. Please try again.')

    return render_template('login.html')


@app.route('/index')
@login_required
def index():
    '''Function to render the dashboard page after successful login.'''
    # Dashboard page after successful login
    return render_template('index.html')

@app.route('/logout')
def logout():
    '''Function to log out the user.'''
    logout_user()
    return redirect(url_for('home'))


# SETINGS PAGE
@app.route('/settings.html')
@login_required  # Ensure that the user is logged in
def settings_html():
    '''Function to render the settings page.'''
    user_id = current_user.get_id()
    user = users_collection.find_one({"_id": ObjectId(user_id)})
    if user:

        return render_template('settings.html', user=user)
    else:
        flash("User not found.")
        return redirect(url_for('index'))
 
#ACCOUNT SETTINGS - UPDATE (FIRST / LAST NAME, EMAIL)
@app.route('/update_account_settings', methods=['POST'])
@login_required
def update_account_settings():
    '''Function to update the user's account settings.'''
    first_name = request.form.get('firstname')
    last_name = request.form.get('lastname')
    email = request.form.get('email')
    

    user_id = current_user.get_id()
    user = users_collection.find_one({"_id": ObjectId(user_id)})

    if user:
        updates = {}
        if user.get('first_name') != first_name:
            updates['first_name'] = first_name

        if user.get('last_name') != last_name:
            updates['last_name'] = last_name

        if user.get('email') != email:
            updates['email'] = email

        if updates:
            users_collection.update_one({"_id":ObjectId(user_id)}, {"$set":updates})
            flash('Your account has been updated successfully')
        else:
            flash('No changes were made to your account.')

    else:
        flash('User not found.')
        return redirect(url_for('settings_html'))

    return redirect(url_for('settings_html'))

@app.route('/change_password', methods=['POST'])
@login_required
def change_password():
    '''Function to change the user's password.'''
    current_password = request.form.get('current-password')
    new_password = request.form.get('new-password')
    confirm_new_password = request.form.get('confirm-new-password')

    user_id = current_user.get_id()
    user = users_collection.find_one({"_id": ObjectId(user_id)})

    if not user:
        return redirect(url_for('settings_html')) #redirect to settings page
    
    #function provided by flask, checks if a plain text password matches a hash password
    #user['password'] is the hashed PW fro DB, current password is what the user enters
    if not bcrypt.check_password_hash(user['password'], current_password): 
        flash('Current password is incorrect.')
        return redirect(url_for('settings_html'))

    if new_password != confirm_new_password:
        flash('New passwords do not match.')
        return redirect(url_for('settings_html'))

    hashed_new_password = bcrypt.generate_password_hash(new_password).decode('utf-8')
    users_collection.update_one({"_id": ObjectId(user_id)}, {"$set": {"password": hashed_new_password}})

    flash('Your password has been updated successfully.')
    return redirect(url_for('settings_html'))

@app.route('/share_via_email', methods=['POST'])
def share_via_email():
    '''Function to share a file via email.'''
    # Get the recipient's email address and the file id from the form data
    recipient = request.form.get('email')
    file_id = request.form.get('file_id')

    # Generate the shareable URL
    share_url = url_for('share_file', file_id=file_id, _external=True)

    # Create a new email message
    msg = Message('Your Shared File',
                  sender='coursify@outlook.com',
                  recipients=[recipient])

    # Add the shareable URL to the email body
    msg.body = f'Here is the file you requested: {share_url}'

    # Send the email
    mail.send(msg)

    return 'Email sent!'

        
@app.route('/share/<file_id>')
def share_file(file_id):
    '''Function to share a file with a unique URL.'''
    # converts id to objectid
    file_id = ObjectId(file_id)

    # Retrieve the file from GridFS
    file = fs.get(file_id)

    # Create a response with the file data
    response = make_response(file.read())
    response.mimetype = 'application/pdf'

    # Set the Content-Disposition header to make the file downloadable
    response.headers.set('Content-Disposition', 'attachment', filename=file.filename)

    return response


def extract_text_from_pdf(pdf_path):
    '''Extract text from a PDF using PdfReader.'''
    # This function extracts text from a PDF using PdfReader
    with open(pdf_path, 'rb') as f:
        reader = PdfReader(f)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
    return text

@app.route('/content')
def my_content():
 '''Function to display the content page.'''
 if current_user.is_authenticated:
    fs=GridFS(db)
    # Get a list of all files in GridFS
    files = fs.find({"user_id": current_user.get_id()}).sort("_id",-1)

    # Create a list to store file data
    file_data = []

    # Loop through all files
    for file in files:
        # Get the file name and the file id (used to retrieve the file later)
        file_data.append({"filename": file.filename, "_id": str(file._id)})

    # Render a template and pass the file data to it
    return render_template('content.html', file_data=file_data)

# @app.route('/content')
# def content():
#     return render_template('content.html')



@app.route('/ai.html')
def ai_html():
    return render_template('ai.html')

@app.route('/preview.html')
def preview():
    return render_template('preview.html')
@app.route('/faq.html')
def faq():
    return render_template('faq.html')
@app.route('/api/chat', methods=['POST'])
@app.route('/api/chat', methods=['POST'])
def chat():
    '''Function for the chatbot API endpoint.'''
    try:
        user_input = request.json['message']
        print("User input received:", user_input)  # Debugging log

        
        
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "AI, your role is to assist users in navigating and utilizing the features of Courisify.ai \
effectively. You will only answer questions related to this website. When a user inquires \
about generating content, guide them through entering the desired content length, difficulty, \
and topics. If they need to manage their notes, direct them to the 'My Content' section. \
For account-related queries, help users find where to update their details or change their \
password on the 'Settings' page. Remember to always respect user privacy and handle their \
data with care. \
Should users need assistance with features or encounter issues, use the AI assistant chat to \
offer support. Provide clear, friendly guidance, ensuring users feel supported at every step. \
For new users, explain the registration process, and for returning users, assist them with the \
login procedure. Always communicate in a warm, professional manner, creating a positive \
user experience. \
Remember, your goal is to help users feel confident in using the website, ensuring they can \
access and utilize all features without difficulty. Keep your instructions simple, and your tone \
friendly, and be ready to answer a range of questions with patience and expertise. \
Here are step-by-step instructions for navigating the website, designed to assist users in \
finding and utilizing the various features offered by Courisify.ai: \
Homepage Navigation: \
Upon landing on the homepage, you will find a navigation bar at the top with the following \
options: HOME, MY CONTENT, SETTINGS, and AI ASSISTANT. There's also a LOGOUT \
button for signing out of your account. The central part of the homepage welcomes \
educators and provides a brief introduction to the site's purpose, emphasizing the automated \
creation of educational content. \
Generating Content: \
To generate content, locate the Generate Content section on the homepage. \
Input the desired Length of Content in words, select the Difficulty level (from 0 to 3), and \
enter the Main Topics for the content you need. You have the option to upload a PDF file by \
clicking Choose File if you want to provide additional material or reference. Click the \
Generate button to start the content creation process. \
Managing Your Content: \
Click on MY CONTENT in the navigation bar to view all the content and notes you have \
created or uploaded. You will see a list of files, each with a unique identifier and a \
timestamp. Click on any file to open, review, or download it. \
Adjusting Settings: \
To update your settings, click on the SETTINGS tab. Here you can update your account \
details such as First Name, Last Name, Password, and Email. To upload or change your \
profile picture, click Choose File under Profile Picture. \
Using the AI Assistant: \
For additional help or to ask questions, click on AI ASSISTANT in the navigation bar.  \
A chat window will open where you can type your message and receive a response from the  \
AI support system. \
Logging In and Out: \
To log in to the site, enter your Email and Password on the login page, and click the Login \
button. You can select Remember Me to stay logged in for future visits. Once you are done \
using the site, click on the LOGOUT button in the navigation bar to securely exit your \
account. \
Registration for New Users: \
If you are a new user, click on the Register button on the homepage. Fill in the registration \
form with your details, including First Name, Last Name, Email, Password, and Confirm \
Password. Click the Register button to create your new account."},
                {"role": "user", "content": user_input}
            ]
        )
        print("Response from OpenAI:", response)  # Debugging log

        ai_reply = response['choices'][0]['message']['content']
        print("AI Reply:", ai_reply)  # Debugging log

        return jsonify({'reply': ai_reply})
    except Exception as e:
        print("An error occurred:", e)  # Debugging log
        return jsonify({'error': str(e)}), 500

def is_latex(text):
    '''Check if a string contains LaTeX content.'''

    return bool(re.search(r"\$.*\$", text))

def render_latex(formula, fontsize=12, dpi=150):
    '''Render a LaTeX formula to an image and return the image as a BytesIO buffer.'''

    # Configure Matplotlib to use LaTeX for rendering
    plt.rcParams['text.usetex'] = True
    
    # Set up a Matplotlib figure and text
    fig = plt.figure()
    fig.patch.set_alpha(0)
    ax = fig.add_subplot(111)
    ax.axis('off')
    ax.patch.set_alpha(0)
    # Use the formula within a TeX environment
    ax.text(0, 0, f"\\begin{{equation}}{formula}\\end{{equation}}", fontsize=fontsize)
    
    # Save the figure to a BytesIO buffer
    buffer = BytesIO()
    fig.savefig(buffer, dpi=dpi, transparent=True, format='png')
    plt.close(fig)
    buffer.seek(0)
    return buffer


def sanitize_filename(text):
    '''Sanitize a string to be used as a filename.'''
    # Remove invalid filename characters
    valid_chars = "-_.() %s%s" % (string.ascii_letters, string.digits)
    sanitized = ''.join(c for c in text if c in valid_chars)
    # Truncate the filename if it's too long
    return sanitized[:255]


def wrap_text(text, max_width, font_name, font_size):
    '''Wrap text to fit within a given width.'''
    # Function to wrap text
    wrapped_text = []
    words = text.split()

    while words:
        line = ''
        while words and stringWidth(line + words[0], font_name, font_size) <= max_width:
            line += (words.pop(0) + ' ')
        wrapped_text.append(line)

    return wrapped_text
def content(prompt, length):
    '''Calls the OpenAI API to generate content based on the user's input.'''
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a content generation tool. Give as much accurate information as you are supposed to giv e and follow the instruction in the prompt."},
                {"role": "user", "content": prompt}
            ]
        )
        return response['choices'][0]['message']['content']
    except Exception as e:
        print(f"Error calling OpenAI API: {e}")
        return None
def call_openai_api(prompt):
    ''' Calls the OpenAI API to generate content based on the user's input.'''
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a content generation tool. Give as much accurate information as you are supposed to giv e and follow the instruction in the prompt."},
                {"role": "user", "content": prompt}
            ]
        )
        return response['choices'][0]['message']['content']
    except Exception as e:
        print(f"Error calling OpenAI API: {e}")
        return None
    

@app.route('/generate_content', methods=['POST'])
def generate_content():
    # Common form data processing
    length = request.form.get('length', type=int, default=100)
    prompt = request.form.get('topics', default='')
    difficulty = request.form.get('difficulty', type=int, default=1)
    content_format = request.form.get('format', default='pdf')

    if content_format == 'pdf':
        return generate_pdf(prompt, length, difficulty)
    elif content_format == 'slides':
        return generate_slides(prompt, length, difficulty)
    else:
        return jsonify(success=False, error="Invalid format selected")

def generate_pdf(prompt, length, difficulty):
    '''Generate a PDF based on the user's input.'''
    # Create a subdirectory for PDFs if it doesn't exist
    pdf_directory = os.path.join(os.getcwd(), 'pdfs')
    if not os.path.exists(pdf_directory):
        os.makedirs(pdf_directory)

    # Process form data
    

    pdf_basename = sanitize_filename(prompt)

    # Ensure the filename is not empty
    pdf_basename = pdf_basename if pdf_basename else 'generated_file'

    # Add a timestamp or random string to the filename to ensure uniqueness
    timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
    unique_suffix = ''.join(random.choices(string.ascii_letters + string.digits, k=6))
    pdf_filename = f"{pdf_basename}_{timestamp}_{unique_suffix}.pdf"

    if difficulty == 1:
        diff = "basic"
    elif difficulty == 2:
        diff = "intermediate"
    elif difficulty == 3:
        diff = "advanced"

    
    toc = call_openai_api("Give Table of contents for topic: " +prompt+"(such that there are max 5 topics and each topic has two sub topics. Topics should be upper Case and subtopics otherwise. Dont Start with heading : Table of contents, just show contents with difficulty " + diff)
    if toc is None:
        return jsonify(success=False, error="Failed to generate text from OpenAI API")
    

    



    
    # Define font properties and margins
    font_name = "Helvetica"
    font_size = 12
    left_margin = 72
    top_margin = 720
    line_height = 14
    page_width, page_height = letter
    bottom_margin = 72
    content_width = page_width - 2 * left_margin

    # Generate PDF
    pdf_path = os.path.join(pdf_directory, pdf_filename)
    c = canvas.Canvas(pdf_path, pagesize=letter)
    c.setFont(font_name, font_size)
    toc_dict = {}
    current_topic = ""

    

    y_position = top_margin
    c.drawString(left_margin, y_position, "Table of Contents")
    y_position -= line_height
    for line in toc.split('\n'):
        
        
        is_topic = line.isupper()  # Assuming topics are in uppercase
        contains_latex = is_latex(line)  # Renamed the variable here

        if contains_latex:
            # Process LaTeX content
            latex_image_io = render_latex(line[1:-1])  # Remove $ symbols
            c.drawImage(latex_image_io, left_margin, y_position, width=200, height=100, preserveAspectRatio=True, anchor='n')
            y_position -= 100  # Adjust for image height
        else:
            # Format as a main topic or subtopic
            if is_topic:
                
                c.setFont("Helvetica-Bold", 14)
                y_position -= 20  # Extra space before a main topic
            else:
                if current_topic:  # Ensure there is a current topic
                    toc_dict[current_topic].append(line.strip())
                c.setFont("Helvetica", 12)
                line = "   " + line  # Indent subtopics

            # Add line to PDF
            c.drawString(left_margin, y_position, line)
            y_position -= line_height

        if y_position < 100:  # Check for end of page and create a new one if needed
            c.showPage()
            c.setFont("Helvetica", 12)
            y_position = top_margin
    
    c.showPage()
    c.setFont("Helvetica", 12)
    y_position = top_margin

    
    lines = toc.split('\n')
    
    print(lines)
    for line in lines:
        is_topic = line.isupper()  # Assuming topics are in uppercase
        contains_latex = is_latex(line)  # Renamed the variable here

        if contains_latex:
            # Process LaTeX content
            latex_image_io = render_latex(line[1:-1])  # Remove $ symbols
            c.drawImage(latex_image_io, left_margin, y_position, width=200, height=100, preserveAspectRatio=True, anchor='n')
            y_position -= 100  # Adjust for image height
            
        else:
            # Format as a main topic or subtopic
            if is_topic:
                c.setFont("Helvetica-Bold", 14)
                wrapped_text = wrap_text(line, content_width, font_name, 14)
                print(wrapped_text)
                
            else:
                
                c.setFont("Helvetica", 12)
                c.drawString(left_margin, y_position, line)
                y_position -= line_height 
                prompt2 = content("explain in 3 lines about the following topic " + line + " related to " + prompt, "length : " + str(length/(len(lines))) + " and difficulty level is " + diff)

                if prompt2 is None:
                    return jsonify(success=False, error="Failed to generate text from OpenAI API")
                wrapped_text = wrap_text(prompt2, content_width, font_name, 12)
                print ("sub topic done")
                
            for text_line in wrapped_text:
                if y_position < bottom_margin:  # Check if we need a new page
                    c.showPage()
                    c.setFont(font_name, 12)
                    y_position = top_margin
                c.drawString(left_margin, y_position, text_line)
                y_position -= line_height 

    # Save the PDF
    c.save()
    print("PDF saved to", pdf_path)



    if current_user.is_authenticated:
        fs = GridFS(db)
        with open(pdf_path, 'rb') as pdf_file:
         fs.put(pdf_file, filename=pdf_filename, user_id=current_user.user_id)
    # 
    #     fs.put(pdf_file, filename=pdf_filename)
    

    # Respond with the URL of the PDF
    pdf_url = url_for('get_pdf', filename=pdf_filename)
    return jsonify(success=True, pdf_url=pdf_url)
def generate_slides(prompt, length, difficulty,):
    '''Generate a presentation based on the user's input.'''
    # Directory where the presentations will be saved
    pptx_directory = os.path.join(os.getcwd(), 'presentations')
    if not os.path.exists(pptx_directory):
        os.makedirs(pptx_directory)


    if difficulty == 1:
        diff = "basic"
    elif difficulty == 2:
        diff = "intermediate"
    elif difficulty == 3:
        diff = "advanced"

    # Define filename for the presentation
    pptx_basename = sanitize_filename(prompt)
    pptx_basename = pptx_basename if pptx_basename else 'generated_presentation'
    timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
    unique_suffix = ''.join(random.choices(string.ascii_letters + string.digits, k=6))
    pptx_filename = f"{pptx_basename}_{timestamp}_{unique_suffix}.pptx"

    # Generate table of contents
    toc = call_openai_api("Give Table of contents for topic: " + prompt + "..." + "Difficulty level is " + diff + "..." + "Length of content is " + str(length) + " words.")
    if toc is None:
        return jsonify(success=False, error="Failed to generate text from OpenAI API")

    # Create a presentation object
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Using a title and content layout

    # Add a slide for the Table of Contents
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Table of Contents"

    content = slide.placeholders[1]
    tf = content.text_frame

    for line in toc.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.level = 0 if line.isupper() else 1

    # Generate and add content for each section
    for line in toc.split('\n'):
        if line.strip():  # Check if line is not empty
            # Generate content for the section
            section_content = call_openai_api(f"Explain {line.strip()} in detail.")

            # Add a new slide for the section
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = line.strip()

            # Add generated content to the slide
            content_box = slide.placeholders[1]
            tf = content_box.text_frame
            tf.text = section_content  # This sets the initial paragraph
            # For more complex formatting, you can add more paragraphs or format this one
            print("Topic done")

    pptx_path = os.path.join(pptx_directory, pptx_filename)
    prs.save(pptx_path)



    if current_user.is_authenticated:
        fs = GridFS(db)
        with open(pptx_path, 'rb') as pptx_file:
            fs.put(pptx_file, filename=pptx_filename, user_id=current_user.user_id)

   

    # Respond with the URL of the presentation
    pptx_url = url_for('get_presentation', filename=pptx_filename)
    return jsonify(success=True, pptx_url=pptx_url)




@app.route('/get_user_pdfs', methods=['GET'])
def get_user_pdfs():
    '''Get all PDFs uploaded by the currently logged in user.'''
    if current_user.is_authenticated:
        fs = GridFS(db)
        user_pdfs = fs.find({'user_id': current_user.user_id})

    return user_pdfs

@app.route('/pdf/<filename>')
def get_pdf(filename):
       '''Download a PDF from GridFS based on its filename.'''
    #  directory = os.path.join(os.getcwd(), 'pdfs')   
    #  file_path = os.path.join(directory, filename)
    #  if os.path.exists(file_path):
    #      return send_from_directory(directory, filename)
    #  else:
    #      return "File not found", 404
       if current_user.is_authenticated:
        fs = GridFS(db)
        grid_out = fs.find_one({'filename': filename, 'user_id': current_user.user_id})

        if grid_out:
            response = make_response(grid_out.read())
            response.mimetype = 'application/pdf'
            return response
        else:
            return "File not found", 404
       else:
        return "Unauthorized", 401
@app.route('/presentation/<filename>')
def get_presentation(filename):
    '''Download a presentation from GridFS based on its filename.'''
    if current_user.is_authenticated:
        fs = GridFS(db)  # Assuming 'db' is your MongoDB database instance
        grid_out = fs.find_one({'filename': filename, 'user_id': current_user.user_id})

        if grid_out:
            # Read the file data from GridFS
            response = make_response(grid_out.read())
            # Set the MIME type to 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            response.mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            # Set the filename in the Content-Disposition header for downloading
            response.headers["Content-Disposition"] = f"attachment; filename={filename}"
            return response
        else:
            # File not found in the database
            abort(404, description="Presentation file not found")
    else:
        # Unauthorized access
        abort(401, description="Unauthorized to access this presentation")

@app.route('/get_doc/<file_id>')
def get_doc(file_id):
    '''Download a file from GridFS based on its file_id.'''
    try:
        file_id = ObjectId(file_id)  # Ensure file_id is a valid ObjectId
        grid_out = fs.get(file_id)
        return send_file(grid_out, attachment_filename=grid_out.filename, as_attachment=True)
    except NoFile:
        return jsonify({'error': 'File not found'}), 404

@app.route('/list_pdfs', methods=['GET'])
def list_pdfs():
    '''List all PDFs in the database.'''
    if current_user.is_authenticated:
        fs = GridFS(db)
        user_pdfs = fs.find({'user_id': current_user.user_id})
        # Now user_pdfs contains only the PDFs generated by the currently logged in user.
        # You can return these PDFs or render them in a template.

    return render_template('content.html', pdfs=user_pdfs)    
  
  #can be used to check if file is in database  
@app.route('/check_file/<filename>', methods=['GET'])
def check_file(filename):
    '''Check if a file exists in the database.'''
    try:
        file = fs.get_last_version(filename=filename)
        if file:
            return jsonify(success=True, message="File exists in the database.")
    except:
        return jsonify(success=False, message="File does not exist in the database.")
    
@app.route('/submit_review', methods=['POST'])
@login_required
def submit_review():
    '''Submit a review for a file.'''
    star_rating = request.form['star_rating']
    review_text = request.form['review_text']
    user_id = ObjectId(current_user.get_id())
    subject = request.form['subject']
    file = request.form['file']
    
    user_details = users_collection.find_one({"_id": user_id})
    
    if user_details is not None:
        first_name = user_details.get('first_name')
        last_name = user_details.get('last_name')

    review = {"user_id": current_user.get_id(), "first_name": first_name, "last_name": last_name, "star_rating": star_rating, "review_text": review_text, "subject": subject, "file": file, "timestamp": datetime.utcnow()}
    reviews_collection.insert_one(review)

    flash('Review submitted successfully.')
    return 'Review sent'
    
@app.route('/reviews')
@login_required
def reviews():
    # Define user_id here
    user_id = ObjectId(current_user.get_id())

    all_reviews = reviews_collection.find().sort("timestamp", -1)
    
    fs = GridFS(db)
    
    user_files = fs.find({"user_id": user_id})
    
    return render_template('reviews.html', reviews=all_reviews, files=user_files)




@app.route('/delete/<file_id>', methods=['POST'])
def delete_file(file_id):
    '''Delete a file from GridFS based on its file_id.'''
    # Convert the file_id to an ObjectId
    file_id = ObjectId(file_id)

    # Delete the file from GridFS
    fs.delete(file_id)

    # Redirect the user to the content page
    return redirect(url_for('my_content'))
 
   

@app.route("/chatbot", methods=["POST"])
def chatbot():
    '''Chatbot route to handle user messages and generate responses.'''
    # Get the message from the POST request
    message = request.json.get("message")

    # Set OpenAI API key (if not already set globally)
    openai.api_key = 'your-api-key'

    # Send the message to OpenAI's API and receive the response
    completion = openai.Completion.create(
        engine="text-davinci-003",  # Use the text-davinci-003 model
        prompt=message,
        max_tokens=150  # Adjust max_tokens if necessary
    )

    if completion.choices and completion.choices[0].text.strip() != "":
        return completion.choices[0].text.strip()
    else:
        return 'Failed to Generate response!'
    
   # User Loader
@login_manager.user_loader
def load_user(user_id):
    '''Load a user from the database using the user_id.'''
    user = users_collection.find_one({"_id": ObjectId(user_id)})
    if not user:
        return None
    return User(user_id=user["_id"], email=user["email"])

@app.route('/generate_quiz/<file_id>')
def generate_quiz(file_id):

    '''Generate a quiz based on the text extracted from a PDF or PPTX file.'''

    # Convert string file_id to ObjectId for MongoDB
    try:
        file_id = ObjectId(file_id)
    except:
        return jsonify({'error': 'Invalid file ID'}), 400

    # Retrieve the file from GridFS
    try:
        file = fs.get(file_id)
    except:
        return jsonify({'error': 'File not found'}), 404


    #Get Name of the file to be used to save the quiz
    file_name = file.filename
    
    # Read the file's contents into a BytesIO stream
    file_stream = io.BytesIO(file.read())

    # Extract text from the PDF or pptx
    file_extension = os.path.splitext(file.filename)[1]
    if file_extension == '.pdf':
        extracted_text = extract_text_from_pdf(file_stream)
    elif file_extension == '.pptx':
        extracted_text = extract_text_from_pptx(file_stream)
    else:
        return jsonify({'error': 'Unsupported file format'}), 400


    # Generate quiz questions based on the extracted text
    questions = generate_questions(extracted_text)

    # Create a Word document for the quiz
    doc = Document()
    doc.add_heading('Quiz', level=1)
    for i, question in enumerate(questions, start=1):
         doc.add_paragraph(f"Q{i}: {question}")  # Corrected line

    # Save the document to a temporary file
    temp_dir = tempfile.mkdtemp()

    timestamp = str(int(time.time()))
    random_string = ''.join(random.choices(string.ascii_lowercase + string.digits, k=8))

    doc_filename = f'quiz_{file_name}.docx'

    
    temp_path = os.path.join(temp_dir, doc_filename)
    doc.save(temp_path)

    # Upload the document to MongoDB using GridFS
    with open(temp_path, 'rb') as doc_file:
        if current_user.is_authenticated:
            file_id = fs.put(doc_file, filename=doc_filename, user_id=str(current_user.get_id()))
    
    # Clean up the temporary file and directory
    os.remove(temp_path)
    os.rmdir(temp_dir)

    # Create a URL for accessing the generated quiz document
    quiz_url = url_for('get_doc', file_id=file_id, _external=True)

    # Respond with the URL of the document
    return jsonify(success=True, quiz_url=quiz_url)

def get_file_stream(file_id):
    '''Retrieve a file from GridFS and return it as a stream.'''
    try:
        file_id = ObjectId(file_id)
        fs = GridFS(db)  # Ensure 'db' is your MongoDB database instance
        file = fs.get(file_id)
        file_stream = io.BytesIO(file.read())
        return file_stream
    except Exception as e:
        print(f"Error retrieving file: {e}")
        return None

def extract_text_from_pdf(file_stream):
    '''Extract text from a PDF file stream.'''
    # Ensure the stream position is at the start
    file_stream.seek(0)
    reader = PdfReader(file_stream)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text
def extract_text_from_pptx(file_stream):
    '''Extract text from a PPTX file stream. It generates questions from table of content which is in first slide of pptx'''
    # Ensure the stream position is at the start
    file_stream.seek(0)
    presentation = Presentation(file_stream)
    text = ""
    # Extracting text from the first slide only
    first_slide = presentation.slides[0]
    for shape in first_slide.shapes:
        if hasattr(shape, "text"):
            text += shape.text + "\n"
            break  # Only extract text from the first question
    return text


def generate_questions(text):
    '''Generate quiz questions based on the given text.'''
    # Construct a prompt for OpenAI API to generate quiz questions
    prompt = f"Create 5 multiple-choice questions based on the following text: \n{text}\nEach question should have 4 options (A, B, C, D), and indicate the correct answer."
    
    # Call the OpenAI API function
    questions_text = call_openai_api(prompt)
    
    if questions_text:
        # Assuming the response text contains the questions formatted as needed
        # You might need to further process this text depending on the response format
        questions = questions_text.strip().split('\n\n')  # Example of simple parsing
        return questions
    else:
        return ["Failed to generate questions."]

    
if __name__ == '__main__':
    app.debug = True
    app.run()

    
@app.route('/submit_review', methods=['POST'], endpoint='submit_review1')
@login_required
def submit_review():
    star_rating = request.form['star_rating']
    review_text = request.form['review_text']
    user_id = ObjectId(current_user.get_id())
    subject = request.form['subject']
   
    review = {
        "user_id": current_user.get_id(),
        "star_rating": star_rating,
        "review_text": review_text,
        "subject": subject,
        "timestamp": datetime.utcnow()  # Optional, for sorting purposes
    }
    reviews_collection.insert_one(review)
   
    flash('Review submitted successfully.')
    return redirect(url_for('reviews'))


@app.route('/reviews', endpoint='reviews1')
@login_required
def reviews():
    all_reviews = reviews_collection.find().sort("timestamp", -1)  # Assuming you want the newest first
    return render_template('reviews.html', reviews=all_reviews)
# function for converting pptx to images

@app.route('/presentation/<filename>')

def pptx_images(pptx_filename):
    '''Converts a PowerPoint (PPTX) file to a PDF and stores it in the pdfs directory.'''

    # Create a subdirectory for PDFs if it doesn't exist
    pdf_directory = os.path.join(os.getcwd(), 'pdfs')
    if not os.path.exists(pdf_directory):
        os.makedirs(pdf_directory)

    # Ensure the filename is safe and generate a unique PDF filename
    base_filename = os.path.splitext(os.path.basename(pptx_filename))[0]
    timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
    unique_suffix = ''.join(random.choices(string.ascii_letters + string.digits, k=6))
    pdf_filename = f"{base_filename}_{timestamp}_{unique_suffix}.pdf"

    # Load the PowerPoint file
    pptx_path = os.path.join(os.getcwd(), pptx_filename)
    presentation = Presentation(pptx_path)

    # Generate PDF
    pdf_path = os.path.join(pdf_directory, pdf_filename)
    c = canvas.Canvas(pdf_path, pagesize=letter)

    # Define basic parameters for the PDF
    font_name = "Helvetica"
    font_size = 12
    left_margin = 72
    top_margin = 750  # Start near the top of the page
    line_height = 14

    c.setFont(font_name, font_size)
    y_position = top_margin

    # Iterate through each slide and then each content in the slide
    for slide_number, slide in enumerate(presentation.slides):
        y_position -= line_height  # Space before new slide content
        c.drawString(left_margin, y_position, f"Slide {slide_number + 1}")
        y_position -= line_height
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    # Simple method to avoid text going beyond the page length
                    if y_position <= 100:  # Near the bottom of the page
                        c.showPage()
                        y_position = top_margin
                        c.setFont(font_name, font_size)
                    c.drawString(left_margin, y_position, text)
                    y_position -= line_height

    # Save the PDF
    c.save()
    print(f"PDF saved to {pdf_path}")

    # Return the path of the generated PDF
    return pdf_path



def pptx_to_images(pptx_file):
    presentation = Presentation(pptx_file)
    images = []
    for i, slide in enumerate(presentation.slides):
        slide_image = slide.export(io.BytesIO())
        images.append(Image.open(slide_image))
    return images
        

if __name__ == '__main__':
    app.run(debug=True)


def create_app():
    app = Flask(__name__)
    return app

